import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# 1. Inicializar Apresentação com formato 16:9 (Widescreen)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Definir Paleta de Cores (Estética Premium Eco-Tech)
COLOR_BG = RGBColor(248, 249, 250)        # Off-white / Cinza Claro
COLOR_CARD_BG = RGBColor(255, 255, 255)   # Branco Puro
COLOR_SLATE = RGBColor(33, 37, 41)         # Slate Escuro (Texto Principal)
COLOR_FOREST = RGBColor(27, 67, 50)        # Verde Floresta (Títulos Principais)
COLOR_EMERALD = RGBColor(45, 106, 79)      # Verde Esmeralda (Destaques e Métricas)
COLOR_MUTED = RGBColor(108, 117, 125)      # Cinza Muted (Subtítulos)
COLOR_ALT_ROW = RGBColor(240, 244, 241)    # Verde muito claro para linhas alternadas

# Estilo de fontes
FONT_TITLE = "Trebuchet MS"
FONT_BODY = "Arial"

# Remover layouts padrões e usar slide em branco para design customizado
blank_layout = prs.slide_layouts[6]

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="PROJETO CIÊNCIAS DO AMBIENTE"):
    # Categoria/Tema pequeno acima do título
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = Inches(0.0)
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = FONT_BODY
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_EMERALD

    # Título do Slide
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = Inches(0.0)
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_FOREST

    # Linha de acento decorativa verde
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_EMERALD
    line.line.color.rgb = COLOR_EMERALD

def create_card(slide, left, top, width, height, title="", border_color=None):
    # Desenhar um retângulo branco (Card) como container estrutural
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.color.rgb = RGBColor(220, 224, 230)
        card.line.width = Pt(1)
        
    # Se houver título no card, adicionar uma caixa de texto dentro dele
    if title:
        tbox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_FOREST

def set_cell_font(cell, text, size=9, bold=False, color=COLOR_SLATE, align=PP_ALIGN.LEFT):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    cell.vertical_anchor = 3 # MSO_ANCHOR_MIDDLE é representado por 3

def style_header_cell(cell, text):
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_FOREST
    set_cell_font(cell, text, size=10, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

def style_body_cell(cell, text, bold=False, color=COLOR_SLATE, bg_color=None, align=PP_ALIGN.LEFT):
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    set_cell_font(cell, text, size=9, bold=bold, color=color, align=align)

# ----------------- SLIDE 1: CAPA -----------------
slide1 = prs.slides.add_slide(blank_layout)
apply_background(slide1)

stripe = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
stripe.fill.solid()
stripe.fill.fore_color.rgb = COLOR_FOREST
stripe.line.fill.background()

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(2.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.0)

p_sub = tf.paragraphs[0]
p_sub.text = "CIÊNCIAS DO AMBIENTE - UTFPR"
p_sub.font.name = FONT_BODY
p_sub.font.size = Pt(12)
p_sub.font.bold = True
p_sub.font.color.rgb = COLOR_EMERALD

p_title = tf.add_paragraph()
p_title.text = "Mobilidade Urbana e Sustentabilidade\nno Centro de Curitiba"
p_title.font.name = FONT_TITLE
p_title.font.size = Pt(44)
p_title.font.bold = True
p_title.font.color.rgb = COLOR_FOREST

p_subtitle = tf.add_paragraph()
p_subtitle.text = "Estrutura do Desenvolvimento e Lógica de Engenharia da Solução"
p_subtitle.font.name = FONT_BODY
p_subtitle.font.size = Pt(18)
p_subtitle.font.color.rgb = COLOR_SLATE

credits_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.2))
tf_cr = credits_box.text_frame
tf_cr.word_wrap = True
tf_cr.margin_left = tf_cr.margin_right = tf_cr.margin_top = tf_cr.margin_bottom = Inches(0.0)

p_cr1 = tf_cr.paragraphs[0]
p_cr1.text = "Equipe Multidisciplinar de Engenharia:"
p_cr1.font.bold = True
p_cr1.font.size = Pt(13)
p_cr1.font.color.rgb = COLOR_SLATE

p_cr2 = tf_cr.add_paragraph()
p_cr2.text = "Giovane Limas Salvi (RA: 2355841) | Rafael Bergo (RA: 2387190) | Gabriel Mororó (RA: 2306298)"
p_cr2.font.size = Pt(12.5)
p_cr2.font.bold = True
p_cr2.font.color.rgb = COLOR_EMERALD

p_cr3 = tf_cr.add_paragraph()
p_cr3.text = "Engenharias de Computação, Elétrica (Automação), e Mecânica"
p_cr3.font.size = Pt(11)
p_cr3.font.color.rgb = COLOR_MUTED


# ----------------- SLIDE 2: INTEGRACAO DAS ENGENHARIAS -----------------
slide2 = prs.slides.add_slide(blank_layout)
apply_background(slide2)
add_header(slide2, "Arquitetura e Integração das Engenharias", "FLUXO GERAL DO TRABALHO")

create_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Engenharia de Sistemas e Fluxo de Dados")
txt_box_left = slide2.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(3.9))
tf_l = txt_box_left.text_frame
tf_l.word_wrap = True
tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = Inches(0.15)

bullets_l2 = [
    "• Abordagem sistêmica e realimentada da mobilidade central.",
    "• Integração de dados lógicos com hardware e eficiência térmica.",
    "• Fluxo dinâmico de tráfego físico gerando estatísticas de filas.",
    "• Algoritmos corretivos reprogramam semáforos no CCO URBS.",
    "• Otimização de torque nos motores reduz CO₂ e poluição acústica.",
    "• Validação dos resultados através de simulação computacional."
]
for idx, text in enumerate(bullets_l2):
    p = tf_l.paragraphs[0] if idx == 0 else tf_l.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(10)

create_card(slide2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Divisão de Responsabilidades Técnicas")
txt_box_right = slide2.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.0), Inches(3.9))
tf_r = txt_box_right.text_frame
tf_r.word_wrap = True
tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = Inches(0.15)

bullets_r2 = [
    "• Computação (Giovane Salvi):",
    "  - Modelagem da malha viária (.osm) e geração probabilística da frota.",
    "  - Desenvolvimento de algoritmos XML de controle atuado (TLS).",
    "• Elétrica & Automação (Rafael Bergo):",
    "  - Projeto de sensoriamento IP virtual e malha de realimentação.",
    "  - Projeto físico-matemático de Onda Verde (Offsets do corredor).",
    "• Engenharia Mecânica (Gabriel Mororó):",
    "  - Dinâmica térmica de aceleração (torque) e curvas HBEFA3.",
    "  - Script Python para parser de logs XML e atenuação acústica."
]
for idx, text in enumerate(bullets_r2):
    p = tf_r.paragraphs[0] if idx == 0 else tf_r.add_paragraph()
    p.text = text
    if "•" in text:
        p.font.bold = True
        p.font.color.rgb = COLOR_FOREST
        p.font.size = Pt(12.5)
        p.space_before = Pt(8)
    else:
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_SLATE
        p.space_after = Pt(4)


# ----------------- SLIDE 3: GIOVANE (COMPUTAÇÃO) -----------------
slide3 = prs.slides.add_slide(blank_layout)
apply_background(slide3)
add_header(slide3, "Computação (Giovane Limas Salvi)", "MODELAGEM E ALGORITMOS DE CONTROLE")

create_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Pipeline de Dados e Distribuição de Frota")
txt_comp_l = slide3.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(1.3))
tf_cl = txt_comp_l.text_frame
tf_cl.word_wrap = True
tf_cl.margin_left = tf_cl.margin_right = tf_cl.margin_top = tf_cl.margin_bottom = Inches(0.0)

bullets_l3 = [
    "• Rede compilada via netconvert com '--junctions.join' para simplificar nós redundantes.",
    "• Script Python automatizando o mapeamento probabilístico da frota central:"
]
for idx, text in enumerate(bullets_l3):
    p = tf_cl.paragraphs[0] if idx == 0 else tf_cl.add_paragraph()
    p.text = text
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(4)

# Adicionar Tabela de Frota no Slide 3 Card 1
table_shape3 = slide3.shapes.add_table(6, 3, Inches(1.1), Inches(3.9), Inches(5.0), Inches(2.2))
table3 = table_shape3.table
table3.columns[0].width = Inches(1.8)
table3.columns[1].width = Inches(1.2)
table3.columns[2].width = Inches(2.0)

style_header_cell(table3.cell(0, 0), "Veículo")
style_header_cell(table3.cell(0, 1), "Proporção")
style_header_cell(table3.cell(0, 2), "Classe no SUMO")

frota_data = [
    ["Carros comuns", "60,0%", "PC_G_EU4 (Gasolina)"],
    ["Motocicletas", "16,0%", "motorcycle (Gasolina)"],
    ["SUVs / Utilitários", "13,0%", "passenger (flex)"],
    ["Vans de Entrega", "5,0%", "delivery (diesel)"],
    ["Caminhões / Ônibus", "6,0%", "HDV_D_EU4 (diesel)"]
]
for r_idx, row in enumerate(frota_data):
    bg_c = COLOR_ALT_ROW if r_idx % 2 == 1 else COLOR_CARD_BG
    style_body_cell(table3.cell(r_idx+1, 0), row[0], bg_color=bg_c)
    style_body_cell(table3.cell(r_idx+1, 1), row[1], bg_color=bg_c, align=PP_ALIGN.CENTER)
    style_body_cell(table3.cell(r_idx+1, 2), row[2], bg_color=bg_c, align=PP_ALIGN.CENTER)

create_card(slide3, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Lógica de Controle Atuado (XML)")
txt_comp_r = slide3.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.0), Inches(3.9))
tf_cr_c = txt_comp_r.text_frame
tf_cr_c.word_wrap = True
tf_cr_c.margin_left = tf_cr_c.margin_right = tf_cr_c.margin_top = tf_cr_c.margin_bottom = Inches(0.15)

bullets_r3 = [
    "• Configuração de semáforos adaptativos (tipo 'actuated' no XML).",
    "• Tempos de fases determinados pela presença real de veículos.",
    "• Parâmetros de extensão de tempo programados na rede:",
    "  - maxGap = 3.0s: limite temporal de espaçamento entre veículos.",
    "  - detectorGap = 2.0s: tempo de extensão de verde por veículo.",
    "  - passingTime = 2.0s: tempo seguro de travessia do detector.",
    "• Execução de roteamento livre de erros viários via duarouter."
]
for idx, text in enumerate(bullets_r3):
    p = tf_cr_c.paragraphs[0] if idx == 0 else tf_cr_c.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(8)


# ----------------- SLIDE 4: RAFAEL (ELÉTRICA) -----------------
slide4 = prs.slides.add_slide(blank_layout)
apply_background(slide4)
add_header(slide4, "Elétrica e Automação (Rafael Bergo)", "SENSORES, CONTROLE E SINCRONISMO")

create_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Arquitetura Física de Sensoriamento IoT")
txt_ele_l = slide4.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(3.9))
tf_el_e = txt_ele_l.text_frame
tf_el_e.word_wrap = True
tf_el_e.margin_left = tf_el_e.margin_right = tf_el_e.margin_top = tf_el_e.margin_bottom = Inches(0.15)

bullets_l4 = [
    "• Laços indutivos virtuais programados na via de tráfego.",
    "• Posicionamento físico dimensionado a 40 metros da retenção.",
    "• Recuo otimizado para conter fila de até 7 automóveis médios.",
    "• Dimensionamento: 7 veículos x 5,5m espaçamento ≈ 38,5m.",
    "• Reaproveitamento de câmeras IP existentes da URBS.",
    "• Regiões de Interesse (ROI) via software eliminam quebra de asfalto."
]
for idx, text in enumerate(bullets_l4):
    p = tf_el_e.paragraphs[0] if idx == 0 else tf_el_e.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(8)

create_card(slide4, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Controle em Malha Fechada e Onda Verde")
txt_ele_r = slide4.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.0), Inches(1.3))
tf_er_e = txt_ele_r.text_frame
tf_er_e.word_wrap = True
tf_er_e.margin_left = tf_er_e.margin_right = tf_er_e.margin_top = tf_er_e.margin_bottom = Inches(0.0)

bullets_r4 = [
    "• Diagrama de blocos de realimentação ativa (Malha Fechada).",
    "• Sincronismo calculado de Onda Verde (Offset = d / v) para velocidade de cruzeiro de 40 km/h (11,11 m/s):"
]
for idx, text in enumerate(bullets_r4):
    p = tf_er_e.paragraphs[0] if idx == 0 else tf_er_e.add_paragraph()
    p.text = text
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(4)

# Adicionar Tabela de Onda Verde no Slide 4 Card 2
table_shape4 = slide4.shapes.add_table(5, 3, Inches(7.1), Inches(3.9), Inches(5.0), Inches(2.2))
table4 = table_shape4.table
table4.columns[0].width = Inches(2.4)
table4.columns[1].width = Inches(1.3)
table4.columns[2].width = Inches(1.3)

style_header_cell(table4.cell(0, 0), "Corredor Floriano")
style_header_cell(table4.cell(0, 1), "Distância")
style_header_cell(table4.cell(0, 2), "Offset")

wave_data = [
    ["Interseção 1 (Referência)", "0 m", "0 s"],
    ["Interseção 2", "222 m", "20 s"],
    ["Interseção 3", "444 m", "40 s"],
    ["Interseção 4", "666 m", "60 s"]
]
for r_idx, row in enumerate(wave_data):
    bg_c = COLOR_ALT_ROW if r_idx % 2 == 1 else COLOR_CARD_BG
    style_body_cell(table4.cell(r_idx+1, 0), row[0], bg_color=bg_c)
    style_body_cell(table4.cell(r_idx+1, 1), row[1], bg_color=bg_c, align=PP_ALIGN.CENTER)
    style_body_cell(table4.cell(r_idx+1, 2), row[2], bg_color=bg_c, align=PP_ALIGN.CENTER)


# ----------------- SLIDE 5: GABRIEL (MECÂNICA) -----------------
slide5 = prs.slides.add_slide(blank_layout)
apply_background(slide5)
add_header(slide5, "Mecânica (Gabriel Mororó)", "TERMODINÂMICA DE COMBUSTÃO E DADOS AMBIENTAIS")

create_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Eficiência Térmica e Equivalência HBEFA3")
txt_mec_l = slide5.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(1.3))
tf_ml = txt_mec_l.text_frame
tf_ml.word_wrap = True
tf_ml.margin_left = tf_ml.margin_right = tf_ml.margin_top = tf_ml.margin_bottom = Inches(0.0)

bullets_l5 = [
    "• Arranques da inércia exigem picos de torque e misturas ricas.",
    "• Integração estequiométrica de curvas HBEFA3 calibradas no SUMO:"
]
for idx, text in enumerate(bullets_l5):
    p = tf_ml.paragraphs[0] if idx == 0 else tf_ml.add_paragraph()
    p.text = text
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(4)

# Adicionar Tabela HBEFA3 no Slide 5 Card 1
table_shape5 = slide5.shapes.add_table(5, 3, Inches(1.1), Inches(3.9), Inches(5.0), Inches(2.2))
table5 = table_shape5.table
table5.columns[0].width = Inches(1.9)
table5.columns[1].width = Inches(1.6)
table5.columns[2].width = Inches(1.5)

style_header_cell(table5.cell(0, 0), "Classe do Veículo")
style_header_cell(table5.cell(0, 1), "Curva HBEFA3")
style_header_cell(table5.cell(0, 2), "CO₂ Idle (Médio)")

emission_data = [
    ["Leves (Gasolina/Flex)", "PC_G_EU4", "2,50 g/min"],
    ["Motos (Gasolina)", "motorcycle", "1,20 g/min"],
    ["Utilitários Médios", "LCV_flex_EU4", "3,20 g/min"],
    ["Pesados (Diesel)", "HDV_D_EU4", "9,50 g/min"]
]
for r_idx, row in enumerate(emission_data):
    bg_c = COLOR_ALT_ROW if r_idx % 2 == 1 else COLOR_CARD_BG
    style_body_cell(table5.cell(r_idx+1, 0), row[0], bg_color=bg_c)
    style_body_cell(table5.cell(r_idx+1, 1), row[1], bg_color=bg_c, align=PP_ALIGN.CENTER)
    style_body_cell(table5.cell(r_idx+1, 2), row[2], bg_color=bg_c, align=PP_ALIGN.CENTER)

create_card(slide5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Parser de Logs XML e Modelagem Acústica")
txt_mec_r = slide5.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.0), Inches(3.9))
tf_mr = txt_mec_r.text_frame
tf_mr.word_wrap = True
tf_mr.margin_left = tf_mr.margin_right = tf_mr.margin_top = tf_mr.margin_bottom = Inches(0.15)

bullets_r5 = [
    "• Desenvolvimento de script Python para ler outputs da simulação.",
    "• Parsing de tripinfo-output.xml e summary.xml gerados na execução.",
    "• Consolidação de dados de tempo de fila, consumo e emissões de CO₂.",
    "• Análise de atenuação acústica urbana baseada nos ciclos de freios:",
    "  - Redução calculada de 3 a 5 dB(A) nos eixos arteriais críticos.",
    "  - Atenua o freia-e-arranca ruidoso de ônibus diesel pesados.",
    "  - Adequação aos limites saudáveis da norma ABNT NBR 10151."
]
for idx, text in enumerate(bullets_r5):
    p = tf_mr.paragraphs[0] if idx == 0 else tf_mr.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(8)


# ----------------- SLIDE 6: RESULTADOS SIMULAÇÃO COMPLETA -----------------
slide6 = prs.slides.add_slide(blank_layout)
apply_background(slide6)
add_header(slide6, "Resultados Práticos: Simulação Completa")

# Adicionar Tabela Grande no Slide 6
rows_r, cols_r = 6, 6
left_r, top_r, width_r, height_r = Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.0)
table_r_shape = slide6.shapes.add_table(rows_r, cols_r, left_r, top_r, width_r, height_r)
table_r = table_r_shape.table

# Definir largura das colunas
table_r.columns[0].width = Inches(3.233)
table_r.columns[1].width = Inches(1.7)
table_r.columns[2].width = Inches(1.7)
table_r.columns[3].width = Inches(1.7)
table_r.columns[4].width = Inches(1.7)
table_r.columns[5].width = Inches(1.73)

# Estilo Header
headers_r = [
    "Métrica de Desempenho", 
    "Manhã As-Is", 
    "Manhã To-Be", 
    "Tarde As-Is", 
    "Tarde To-Be", 
    "Impacto Máx."
]
for col_idx, text in enumerate(headers_r):
    style_header_cell(table_r.cell(0, col_idx), text)

# Dados
data_r = [
    ["Veículos Concluídos (Escoamento)", "15.913 veíc.", "17.118 veíc.", "16.362 veíc.", "19.191 veíc.", "+17.29% (Vazão)"],
    ["Tempo Médio de Viagem", "703,6 s", "436,4 s", "754,3 s", "455,9 s", "-39.56% (Viagem)"],
    ["Tempo Médio de Espera (Filas)", "314,7 s", "53,1 s", "361,5 s", "67,2 s", "-83.11% (Espera)"],
    ["Consumo de Combustível", "11.343.957 L", "8.208.770 L", "12.448.560 L", "9.633.984 L", "-27.64% (Consumo)"],
    ["Emissões de CO₂", "35.684 kg", "25.834 kg", "39.162 kg", "30.321 kg", "-27.60% (CO₂)"]
]

for row_idx, row_data in enumerate(data_r):
    bg = COLOR_ALT_ROW if row_idx % 2 == 1 else COLOR_CARD_BG
    for col_idx, text in enumerate(row_data):
        cell = table_r.cell(row_idx + 1, col_idx)
        bold = True if col_idx == 0 or col_idx == 5 else False
        color = COLOR_EMERALD if col_idx == 5 else COLOR_SLATE
        align = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
        style_body_cell(cell, text, bold=bold, color=color, bg_color=bg, align=align)

note_box = slide6.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(11.733), Inches(0.9))
tf_n = note_box.text_frame
tf_n.word_wrap = True
tf_n.margin_left = tf_n.margin_right = tf_n.margin_top = tf_n.margin_bottom = Inches(0.1)

bullets_n = [
    "• Lógica Adaptativa no CCO: O tempo de fila caiu de 6 minutos para apenas 1 minuto (redução de mais de 81% na tarde).",
    "• Eficiência de Consumo: Redução de 27% nas emissões globais da rede, salvando 18,69 toneladas de CO₂ por dia."
]
for idx, text in enumerate(bullets_n):
    p = tf_n.paragraphs[0] if idx == 0 else tf_n.add_paragraph()
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_FOREST
    p.space_after = Pt(4)


# ----------------- SLIDE 6B: COMPARATIVO GRÁFICO (NOVO) -----------------
slide6b = prs.slides.add_slide(blank_layout)
apply_background(slide6b)
add_header(slide6b, "Análise Gráfica de Tráfego e Emissões", "COMPARATIVO VISUAL AS-IS VS TO-BE")

# Card Esquerdo: Gráfico de Tempos
create_card(slide6b, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Melhoria nos Tempos de Tráfego (Pico Tarde)")

chart_data_time = CategoryChartData()
chart_data_time.categories = ['Tempo de Viagem (s)', 'Tempo em Fila (s)']
chart_data_time.add_series('As-Is (Estático)', (754.3, 361.5))
chart_data_time.add_series('To-Be (Atuado)', (455.9, 67.2))

x_t, y_t, cx_t, cy_t = Inches(1.1), Inches(2.6), Inches(5.0), Inches(3.7)
chart_time = slide6b.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x_t, y_t, cx_t, cy_t, chart_data_time
).chart
chart_time.has_legend = True
chart_time.legend.position = XL_LEGEND_POSITION.BOTTOM
chart_time.legend.include_in_layout = False

# Card Direito: Gráfico de CO2
create_card(slide6b, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Redução de Emissões de CO₂ (kg)")

chart_data_co2 = CategoryChartData()
chart_data_co2.categories = ['Pico Manhã (9000s)', 'Pico Tarde (7200s)']
chart_data_co2.add_series('As-Is (Estático)', (35684.0, 39162.0))
chart_data_co2.add_series('To-Be (Atuado)', (25834.0, 30321.0))

x_c, y_c, cx_c, cy_c = Inches(7.1), Inches(2.6), Inches(5.0), Inches(3.7)
chart_co2 = slide6b.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x_c, y_c, cx_c, cy_c, chart_data_co2
).chart
chart_co2.has_legend = True
chart_co2.legend.position = XL_LEGEND_POSITION.BOTTOM
chart_co2.legend.include_in_layout = False


# ----------------- SLIDE 7: EFICIÊNCIA ECOLÓGICA -----------------
slide7 = prs.slides.add_slide(blank_layout)
apply_background(slide7)
add_header(slide7, "Eficiência Ecológica e Prova Matemática")

create_card(slide7, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8), "Cálculo da Emissão Evitada em Marcha Lenta")
txt_eco_l = slide7.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.4), Inches(1.4))
tf_el = txt_eco_l.text_frame
tf_el.word_wrap = True
tf_el.margin_left = tf_el.margin_right = tf_el.margin_top = tf_el.margin_bottom = Inches(0.0)

p_el1 = tf_el.paragraphs[0]
p_el1.text = "CO₂ Evitado = N × Δt × Fator_idle"
p_el1.font.bold = True
p_el1.font.size = Pt(18)
p_el1.font.color.rgb = COLOR_EMERALD
p_el1.alignment = PP_ALIGN.CENTER
p_el1.space_after = Pt(6)

p_el2 = tf_el.add_paragraph()
p_el2.text = "• Mapeamento de parâmetros dinâmicos da tarde para alimentar a equação de queima estequiométrica:"
p_el2.font.size = Pt(11.5)
p_el2.font.color.rgb = COLOR_SLATE
p_el2.space_after = Pt(4)

# Adicionar Tabela de Parâmetros de Equação no Slide 7 Card 1
table_shape7 = slide7.shapes.add_table(5, 3, Inches(1.1), Inches(3.9), Inches(5.4), Inches(2.2))
table7 = table_shape7.table
table7.columns[0].width = Inches(1.2)
table7.columns[1].width = Inches(2.7)
table7.columns[2].width = Inches(1.5)

style_header_cell(table7.cell(0, 0), "Parâmetro")
style_header_cell(table7.cell(0, 1), "Descrição Física")
style_header_cell(table7.cell(0, 2), "Valor Aplicado")

calc_data = [
    ["N", "Veículos escoados (tarde)", "19.191 veíc."],
    ["Δt", "Espera ociosa poupada", "4,905 min / veíc."],
    ["Fator_idle", "Emissão média em marcha lenta", "3,12 g/min"],
    ["CO₂ Salvo", "Emissão evitada direta", "293,7 kg CO₂"]
]
for r_idx, row in enumerate(calc_data):
    bg_c = COLOR_ALT_ROW if r_idx % 2 == 1 else COLOR_CARD_BG
    style_body_cell(table7.cell(r_idx+1, 0), row[0], bg_color=bg_c)
    style_body_cell(table7.cell(r_idx+1, 1), row[1], bg_color=bg_c)
    style_body_cell(table7.cell(r_idx+1, 2), row[2], bg_color=bg_c, align=PP_ALIGN.CENTER)

create_card(slide7, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), "Impacto Ecológico Central Consolidado")
txt_eco_r = slide7.shapes.add_textbox(Inches(7.5), Inches(2.5), Inches(4.7), Inches(3.9))
tf_er = txt_eco_r.text_frame
tf_er.word_wrap = True
tf_er.margin_left = tf_er.margin_right = tf_er.margin_top = tf_er.margin_bottom = Inches(0.15)

bullets_r7 = [
    "• Economia Atmosférica Diária de CO₂:",
    "  - Redução diária nos picos: 18,69 toneladas de CO₂.",
    "  - Redução acumulada por ano letivo: mais de 4.800 toneladas.",
    "  - Equivalente ecológico de plantar mais de 34.000 árvores.",
    "• Redução e Melhoria da Pressão Acústica:",
    "  - Queda de 3 a 5 dB(A) medida nos eixos de maior tráfego.",
    "  - Eliminação de picos de ruído mecânico de veículos a diesel.",
    "  - Níveis acústicos adequados aos padrões da NBR 10151."
]
for text in bullets_r7:
    p = tf_er.paragraphs[0] if tf_er.paragraphs[0].text == "" else tf_er.add_paragraph()
    p.text = text
    if "•" in text:
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_FOREST
        p.space_before = Pt(6)
    else:
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_SLATE
        p.space_after = Pt(3)


# ----------------- SLIDE 8: INTEGRACAO OPERACIONAL -----------------
slide8 = prs.slides.add_slide(blank_layout)
apply_background(slide8)
add_header(slide8, "Integração Operacional: CCO URBS & Firmware", "VIABILIDADE FÍSICA E COMUNICACIONAL")

create_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Arquitetura Física da Solução")
txt_via_l = slide8.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(3.9))
tf_vl = txt_via_l.text_frame
tf_vl.word_wrap = True
tf_vl.margin_left = tf_vl.margin_right = tf_vl.margin_top = tf_vl.margin_bottom = Inches(0.15)

bullets_l8 = [
    "• Fibra óptica ligando câmeras e semáforos ao CCO URBS.",
    "• Transmissão de streams de vídeo centralizada em tempo real.",
    "• Servidores centrais de análise de imagens geram laços virtuais.",
    "• Instruções de controle dinâmicas são enviadas aos cruzamentos.",
    "• Aproveitamento de infraestrutura evita escavações de pavimentos."
]
for idx, text in enumerate(bullets_l8):
    p = tf_vl.paragraphs[0] if idx == 0 else tf_vl.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(10)

create_card(slide8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Lógica do Firmware e Protocolo de Controle")
txt_via_r = slide8.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.0), Inches(1.3))
tf_vr = txt_via_r.text_frame
tf_vr.word_wrap = True
tf_vr.margin_left = tf_vr.margin_right = tf_vr.margin_top = tf_vr.margin_bottom = Inches(0.0)

bullets_r8 = [
    "• Comunicação local entre CCO e campo via NTCIP 1202 e sincronismo via NTP/GPS. Detalhes dos componentes:"
]
for idx, text in enumerate(bullets_r8):
    p = tf_vr.paragraphs[0] if idx == 0 else tf_vr.add_paragraph()
    p.text = text
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(4)

# Adicionar Tabela de Componentes no Slide 8 Card 2
table_shape8 = slide8.shapes.add_table(5, 3, Inches(7.1), Inches(3.9), Inches(5.0), Inches(2.2))
table8 = table_shape8.table
table8.columns[0].width = Inches(1.6)
table8.columns[1].width = Inches(1.6)
table8.columns[2].width = Inches(1.8)

style_header_cell(table8.cell(0, 0), "Componente")
style_header_cell(table8.cell(0, 1), "Protocolo / Tec.")
style_header_cell(table8.cell(0, 2), "Função Operacional")

component_data = [
    ["Câmeras IP", "H.264 / ONVIF", "Mapeamento ROI virtual"],
    ["Servidor CCO", "Python / OpenCV", "Cálculo de fila local"],
    ["Semáforos", "NTCIP 1202", "Força fase (Force Off)"],
    ["Sincronismo", "Servidores NTP", "Estabilidade de offsets"]
]
for r_idx, row in enumerate(component_data):
    bg_c = COLOR_ALT_ROW if r_idx % 2 == 1 else COLOR_CARD_BG
    style_body_cell(table8.cell(r_idx+1, 0), row[0], bg_color=bg_c)
    style_body_cell(table8.cell(r_idx+1, 1), row[1], bg_color=bg_c, align=PP_ALIGN.CENTER)
    style_body_cell(table8.cell(r_idx+1, 2), row[2], bg_color=bg_c, align=PP_ALIGN.CENTER)


# ----------------- SLIDE 9: CONCLUSÕES -----------------
slide9_c = prs.slides.add_slide(blank_layout)
apply_background(slide9_c)
add_header(slide9_c, "Conclusões Técnicas e Integração Prática", "SÍNTESE DAS ENGEMHARIAS")

create_card(slide9_c, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8), "Validação da Abordagem Multidisciplinar")
txt_concl = slide9_c.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11.1), Inches(3.9))
tf_co = txt_concl.text_frame
tf_co.word_wrap = True
tf_co.margin_left = tf_co.margin_right = tf_co.margin_top = tf_co.margin_bottom = Inches(0.15)

bullets_l9 = [
    "• Fim do isolamento entre engenharias no design de soluções ambientais.",
    "• Algoritmos lógicos (Computação) conectados ao sincronismo físico (Elétrica) e validados termodinamicamente (Mecânica).",
    "• Otimização semafórica por software como alternativa de alta viabilidade e baixíssimo custo.",
    "• Redução líquida imediata de 18,69 toneladas de CO₂ diárias na malha central de Curitiba.",
    "• A engenharia aplicada diretamente ao bem-estar e à saúde pública no ambiente urbano."
]
for idx, text in enumerate(bullets_l9):
    p = tf_co.paragraphs[0] if idx == 0 else tf_co.add_paragraph()
    p.text = text
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_SLATE
    p.space_after = Pt(10)

# Frase de fechamento
p_quote = tf_co.add_paragraph()
p_quote.text = "\n\"A engenharia multidisciplinar transforma a infraestrutura de tráfego em uma barreira de proteção socioambient.\""
p_quote.font.italic = True
p_quote.font.bold = True
p_quote.font.size = Pt(14)
p_quote.font.color.rgb = COLOR_EMERALD
p_quote.alignment = PP_ALIGN.CENTER


# ----------------- SLIDE 10: PERGUNTAS E RESPOSTAS -----------------
slide10_c = prs.slides.add_slide(blank_layout)
apply_background(slide10_c)

qa_box = slide10_c.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
tf_qa = qa_box.text_frame
tf_qa.word_wrap = True
tf_qa.margin_left = tf_qa.margin_right = tf_qa.margin_top = tf_qa.margin_bottom = Inches(0.1)

p_qa1 = tf_qa.paragraphs[0]
p_qa1.text = "Obrigado!"
p_qa1.font.name = FONT_TITLE
p_qa1.font.size = Pt(54)
p_qa1.font.bold = True
p_qa1.font.color.rgb = COLOR_FOREST
p_qa1.alignment = PP_ALIGN.CENTER

p_qa2 = tf_qa.add_paragraph()
p_qa2.text = "Perguntas & Discussão"
p_qa2.font.name = FONT_BODY
p_qa2.font.size = Pt(24)
p_qa2.font.color.rgb = COLOR_EMERALD
p_qa2.alignment = PP_ALIGN.CENTER

p_qa3 = tf_qa.add_paragraph()
p_qa3.text = "\nGrupo de Consultoria Multidisciplinar UTFPR"
p_qa3.font.name = FONT_BODY
p_qa3.font.size = Pt(14)
p_qa3.font.color.rgb = COLOR_MUTED
p_qa3.alignment = PP_ALIGN.CENTER

# Salvar apresentação
output_path = "Apresentacao_Projeto.pptx"
try:
    prs.save(output_path)
    print(f"[SUCESSO] Apresentação salva em {output_path}")
except PermissionError:
    output_path = "Apresentacao_Projeto_v2.pptx"
    prs.save(output_path)
    print(f"[SUCESSO] O arquivo principal estava aberto em outro programa. Salvo alternativamente em {output_path}")
